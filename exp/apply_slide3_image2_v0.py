from __future__ import print_function

import os
import zipfile

from pptx import Presentation
from pptx.util import Inches


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
SRC_PPTX = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v0.pptx")
OUT_PPTX = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v1_image2.pptx")
IMAGE = os.path.join(ROOT, "assets", "generated", "slide3_image2_v0.png")


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def count_notes(path):
    with zipfile.ZipFile(path, "r") as zf:
        return len([name for name in zf.namelist() if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")])


def main():
    if not os.path.exists(SRC_PPTX):
        raise RuntimeError("source PPTX not found: %s" % SRC_PPTX)
    if not os.path.exists(IMAGE):
        raise RuntimeError("generated image not found: %s" % IMAGE)

    prs = Presentation(SRC_PPTX)
    slide = prs.slides[2]

    # Preserve header, rule, page number, footer, and notes. Replace only the old body framework.
    body_top = Inches(1.55)
    body_bottom = Inches(6.95)
    for shape in list(slide.shapes):
        if shape.top >= body_top and shape.top < body_bottom:
            remove_shape(shape)

    pic_w = Inches(9.65)
    pic_h = Inches(5.43)
    pic_x = int((prs.slide_width - pic_w) / 2)
    pic_y = Inches(1.55)
    pic = slide.shapes.add_picture(IMAGE, pic_x, pic_y, width=pic_w, height=pic_h)
    pic.name = "Slide 3 image2 academic framework diagram"

    prs.save(OUT_PPTX)
    print("pptx:", OUT_PPTX)
    print("slide_count:", len(prs.slides))
    print("notes_slides:", count_notes(OUT_PPTX))
    print("image:", IMAGE)


if __name__ == "__main__":
    main()
