from __future__ import print_function

import os
import zipfile

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
SRC = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v2_template.pptx")
OUT = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v3_template_align.pptx")


def count_notes(path):
    with zipfile.ZipFile(path, "r") as zf:
        return len([n for n in zf.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")])


def main():
    prs = Presentation(SRC)
    slide = prs.slides[1]

    # Body cards on slide 2, identified by position to avoid relying on localized shape names.
    for shape in slide.shapes:
        x = round(shape.left / 914400, 2)
        y = round(shape.top / 914400, 2)
        w = round(shape.width / 914400, 2)
        h = round(shape.height / 914400, 2)
        if y == 2.52 and w == 3.35 and h == 2.65 and hasattr(shape, "text_frame"):
            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(5)
                p.line_spacing = 1.12
                for r in p.runs:
                    r.font.size = Pt(15.5)
                    r.font.name = "Microsoft YaHei"

    prs.save(OUT)
    print("pptx:", OUT)
    print("notes_slides:", count_notes(OUT))


if __name__ == "__main__":
    main()
