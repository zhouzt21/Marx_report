from __future__ import print_function

import os
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
SRC = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v4_researched_slide2.pptx")
OUT = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v5_final_report_slide2.pptx")


TEXT = {
    0: "案例基础：深圳歌尔泰克走访与歌尔公开资料",
    3: "案例研究聚焦高科技企业在全球竞争中的技术能力积累",
    4: "案例对象",
    5: "深圳歌尔泰克科技有限公司\n线下走访、展厅参观\n座谈交流与实践推文",
    7: "产业背景",
    8: "声光电精密零组件\n智能声学整机、智能硬件\nAI 眼镜、MR/AR 等方向",
    10: "分析维度",
    11: "研发设计、工程转化\n产业链协同、人才组织\n全球布局",
    12: "案例意义：从企业走访出发，观察高科技企业如何在细分技术领域和全球产业链环境中提升自身能力。",
    13: "来源：实践推文；歌尔股份 2025 年报/摘要；周文、许凌云 2023；尹西明等 2024；Yin 2018",
}

NOTES = """这一页对应第八部分的案例基础。这里不再把材料来源当作写作说明，而是直接说明案例本身为什么适合分析技术突围。

本次案例以深圳市歌尔泰克科技有限公司的走访为主要基础。我们掌握的材料包括线下参观、展厅讲解、座谈交流和实践推文。结合歌尔股份 2025 年年度报告及摘要，可以看到相关业务集中在声光电精密零组件、智能声学整机和智能硬件等方向，也涉及 AI 智能眼镜、MR/AR 等新兴智能终端。

因此，第二页的重点是说明案例研究将从研发设计、工程转化、产业链协同、人才组织和全球布局五个维度展开。这些维度既来自走访观察，也能同新质生产力、企业创新主体和全球价值链升级等文献形成对应。"""


def count_notes(path):
    with zipfile.ZipFile(path, "r") as zf:
        return len([n for n in zf.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")])


def set_text(shape, text, font_size=None, bold=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, color=(20, 28, 42)):
    shape.text = text
    tf = shape.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for p in tf.paragraphs:
        p.alignment = align
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.color.rgb = RGBColor(*color)
            if font_size is not None:
                r.font.size = Pt(font_size)
            if bold is not None:
                r.font.bold = bold


def update_notes(slide):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = NOTES


def main():
    prs = Presentation(SRC)
    slide = prs.slides[1]

    for idx, shape in enumerate(slide.shapes):
        if idx not in TEXT or not getattr(shape, "has_text_frame", False):
            continue
        if idx == 0:
            set_text(shape, TEXT[idx], font_size=22, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        elif idx == 3:
            set_text(shape, TEXT[idx], font_size=17, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        elif idx in (4, 7, 10):
            set_text(shape, TEXT[idx], font_size=16, bold=True, color=(255, 255, 255))
        elif idx in (5, 8, 11):
            set_text(shape, TEXT[idx], font_size=14.2, bold=False)
        elif idx == 12:
            set_text(shape, TEXT[idx], font_size=11.2, bold=False, color=(55, 55, 70))
        elif idx == 13:
            set_text(shape, TEXT[idx], font_size=8.8, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, color=(32, 32, 32))

    update_notes(slide)
    prs.save(OUT)
    print("pptx:", OUT)
    print("notes_slides:", count_notes(OUT))


if __name__ == "__main__":
    main()

