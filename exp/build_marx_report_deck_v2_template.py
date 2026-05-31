from __future__ import print_function

import datetime as dt
import json
import os
import re
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
ALIGN = os.path.join(ROOT, "align")
TEMPLATE = os.path.join(ROOT, "2025nianduyanshiwengaomoban2-tongyongzhuti.pptx")
LAYOUT_PLAN = os.path.join(ALIGN, "ppt_layout_plan_v0.json")
NOTES_ARTIFACT = os.path.join(ALIGN, "ppt_speaker_notes_rehearsal_v0.md")
IMAGE2 = os.path.join(ROOT, "assets", "generated", "slide3_image2_v0_crop_source.png")
OUT_DIR = os.path.join(ROOT, "generated_pptx_test")
PPTX_OUT = os.path.join(OUT_DIR, "marx_report_7_8_section_v2_template.pptx")
MANIFEST_OUT = os.path.join(ALIGN, "ppt_deck_build_manifest_v2_template.md")

PRIMARY = RGBColor(0x3C, 0x36, 0x52)
PURPLE = RGBColor(0x91, 0x2C, 0x8D)
ACCENT1 = RGBColor(0x75, 0x61, 0xD6)
ACCENT2 = RGBColor(0x51, 0x99, 0xEA)
ACCENT3 = RGBColor(0x68, 0xA4, 0xC6)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFC)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
MID_LINE = RGBColor(0xD8, 0xD9, 0xE6)
GRAY = RGBColor(0x67, 0x69, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def rel(path):
    return os.path.relpath(path, ROOT).replace("\\", "/")


def delete_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    r_id = slide_ids[index].rId
    prs.part.drop_rel(r_id)
    slide_id_list.remove(slide_ids[index])


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def keep_template_header_only(slide, title):
    for shape in list(slide.shapes):
        if shape.top >= Inches(1.05):
            remove_shape(shape)

    title_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.left < Inches(5.5) and shape.top < Inches(0.7):
            title_shape = shape
            break
    if title_shape is not None:
        tf = title_shape.text_frame
        tf.clear()
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(25)
        r.font.bold = True
        r.font.color.rgb = PURPLE


def add_textbox(slide, x, y, w, h, text, font_size=16, bold=False,
                color=PRIMARY, align=PP_ALIGN.LEFT, font="Microsoft YaHei",
                margin=0.05, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def set_shape_text(shape, text, size=16, bold=False, color=PRIMARY,
                   align=PP_ALIGN.CENTER, font="Microsoft YaHei"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_action_title(slide, text):
    add_textbox(slide, 0.55, 1.02, 11.8, 0.5, text, 19, True, PRIMARY)


def add_footer(slide, text, page_no):
    add_textbox(slide, 0.52, 7.06, 10.9, 0.23, text, 8.2, False, GRAY)
    add_textbox(slide, 12.2, 7.06, 0.55, 0.23, "%02d" % page_no, 8.2, False, GRAY, align=PP_ALIGN.RIGHT)


def rounded_box(slide, x, y, w, h, text="", fill=PANEL, line=MID_LINE,
                size=15, bold=False, color=PRIMARY, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(radius_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    if text:
        set_shape_text(box, text, size=size, bold=bold, color=color)
    return box


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT2, width=1.8):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def extract_slide_notes():
    text = open(NOTES_ARTIFACT, "r", encoding="utf-8").read()
    notes = {}
    matches = list(re.finditer(r"^### Slide (\d+)\. .*$", text, flags=re.M))
    for idx, match in enumerate(matches):
        slide_no = int(match.group(1))
        start = match.start()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else text.find("\n## ", match.end())
        if end == -1:
            end = len(text)
        notes[slide_no] = text[start:end].strip()
    return notes


def write_notes(slide, note_text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = note_text


def count_notes(path):
    with zipfile.ZipFile(path, "r") as zf:
        return len([n for n in zf.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")])


def build_deck():
    prs = Presentation(TEMPLATE)
    selected = {10, 12, 22, 25}  # template slides 11, 13, 23, 26
    for index in range(len(prs.slides) - 1, -1, -1):
        if index not in selected:
            delete_slide(prs, index)

    with open(LAYOUT_PLAN, "r", encoding="utf-8") as f:
        plan = json.load(f)
    notes = extract_slide_notes()

    # Slide 1
    s = prs.slides[0]
    keep_template_header_only(s, "研究方法：多源互证闭环")
    add_action_title(s, "第 7-8 部分将研究方法和案例观察连接成一个可验证的分析闭环")
    rounded_box(s, 0.72, 1.78, 11.9, 4.95, "", WHITE, MID_LINE)
    add_textbox(s, 1.0, 2.02, 2.2, 0.35, "方法链条", 17, True, PURPLE)
    nodes = ["文献研究", "案例研究", "实地走访", "比较分析", "归纳分析"]
    xs = [1.05, 3.25, 5.45, 7.65, 9.85]
    y = 3.35
    for i, node in enumerate(nodes):
        color = PURPLE if i in (0, 4) else ACCENT2
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xs[i] + 0.55), Inches(y - 0.78), Inches(0.44), Inches(0.44))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        set_shape_text(badge, str(i + 1), 10.5, True, WHITE)
        rounded_box(s, xs[i], y, 1.55, 0.72, node, LIGHT_BG, color, 14, True, PRIMARY)
        if i < len(nodes) - 1:
            add_arrow(s, xs[i] + 1.62, y + 0.36, xs[i + 1] - 0.08, y + 0.36, color, 1.5)
    add_arrow(s, 10.45, 4.55, 1.45, 4.55, ACCENT3, 1.4)
    add_textbox(s, 4.55, 4.75, 4.2, 0.28, "理论框架 - 案例材料 - 路径归纳", 13.5, True, ACCENT1, align=PP_ALIGN.CENTER)
    rounded_box(s, 2.25, 5.65, 8.95, 0.55,
                "材料不是孤立陈列，而是在方法链条中互相校正、共同支撑后续分析。",
                RGBColor(0xF5, 0xF2, 0xFA), MID_LINE, 14, False)
    add_footer(s, "来源：原报告、实践推文、公开资料、相关文献整理", 1)
    write_notes(s, notes[1])

    # Slide 2
    s = prs.slides[1]
    keep_template_header_only(s, "案例基础：以深圳歌尔泰克走访材料为入口")
    add_action_title(s, "多源互证能避免案例分析停留在参观感想或企业介绍层面")
    columns = [
        ("走访入口", ["深圳歌尔泰克走访", "参观与座谈", "产品展示与成员观察"], PURPLE),
        ("公开资料", ["歌尔股份年报/摘要", "业务与研发背景", "全球布局信息"], ACCENT2),
        ("分析用途", ["案例典型性", "技术突围维度", "后续路径归纳"], ACCENT3),
    ]
    for i, (header, items, color) in enumerate(columns):
        x = 0.78 + i * 4.05
        rounded_box(s, x, 1.9, 3.35, 0.52, header, color, color, 16, True, WHITE)
        body = rounded_box(s, x, 2.52, 3.35, 2.65, "", WHITE, MID_LINE)
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.22)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.3)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(14.5)
            p.font.color.rgb = PRIMARY
            p.space_after = Pt(7)
        if i < 2:
            add_arrow(s, x + 3.44, 3.82, x + 3.85, 3.82, color, 1.8)
    rounded_box(s, 2.85, 5.72, 7.65, 0.62,
                "口径：以深圳歌尔泰克走访材料为入口，结合歌尔股份公开资料补充背景。",
                RGBColor(0xF5, 0xF2, 0xFA), ACCENT1, 14.5, True)
    add_footer(s, "来源：实践推文；歌尔股份公开资料；相关文献整理", 2)
    write_notes(s, notes[2])

    # Slide 3
    s = prs.slides[2]
    keep_template_header_only(s, "初步观察：技术突围的五个维度")
    add_action_title(s, "歌尔走访材料显示技术突围依赖研发、工程、供应链、人才和全球布局的系统协同")
    rounded_box(s, 1.38, 1.52, 10.6, 5.42, "", WHITE, MID_LINE)
    s.shapes.add_picture(IMAGE2, Inches(1.92), Inches(1.55), width=Inches(9.5), height=Inches(5.34))
    add_footer(s, "来源：实践推文；公开资料；新质生产力、企业创新主体与全球价值链相关文献", 3)
    write_notes(s, notes[3])

    # Slide 4
    s = prs.slides[3]
    keep_template_header_only(s, "参考文献与资料来源")
    add_action_title(s, "参考文献为方法、理论和案例材料提供最低限度支撑")
    refs = plan["deck"]["slides"][3]["content_structure"][1]["rows"]
    table_shape = s.shapes.add_table(len(refs) + 1, 2, Inches(0.72), Inches(1.72), Inches(11.9), Inches(4.85))
    table = table_shape.table
    table.columns[0].width = Inches(8.2)
    table.columns[1].width = Inches(3.7)
    headers = ["文献/资料", "支撑点"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for r, (ref_text, support) in enumerate(refs, start=1):
        table.cell(r, 0).text = ref_text
        table.cell(r, 1).text = support
        for c in range(2):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF8, 0xF9, 0xFC)
    for r in range(len(refs) + 1):
        for c in range(2):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(10.2 if r else 12)
                p.font.bold = bool(r == 0)
                p.font.color.rgb = WHITE if r == 0 else PRIMARY
    add_textbox(s, 0.75, 6.72, 11.3, 0.24,
                "注：2025 年报条目需在最终合并前按组内统一格式和实际题录核对。",
                8.4, False, GRAY)
    add_footer(s, "来源：实践推文；公开资料；相关文献整理", 4)
    write_notes(s, notes[4])

    return prs


def collect_text(prs):
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def validate(path):
    prs = Presentation(path)
    full_text = collect_text(prs)
    checks = {
        "file_exists": os.path.exists(path),
        "opened_by_python_pptx": True,
        "slide_count": len(prs.slides),
        "aspect_ratio_16_9": prs.slide_width == 12192000 and prs.slide_height == 6858000,
        "template_header_logo_present": all(any(s.shape_type == 13 and s.top < Inches(0.8) and s.left > Inches(10.5) for s in slide.shapes) for slide in prs.slides),
        "slide_2_entity_boundary": "以深圳歌尔泰克走访材料为入口，结合歌尔股份公开资料补充背景" in full_text,
        "slide_3_initial_observation": "初步观察" in full_text,
        "forbidden_wording_absent": "已证明" not in full_text and "最终模型" not in full_text,
        "placeholder_text_absent": not any(s in full_text for s in ["添加标题", "添加正文", "输入图片", "XXX", "关键词"]),
        "notes_slides_present": True,
    }
    try:
        checks["notes_chars"] = [len(slide.notes_slide.notes_text_frame.text) for slide in prs.slides]
        checks["notes_slides_present"] = all(v > 0 for v in checks["notes_chars"])
    except Exception:
        checks["notes_slides_present"] = False
    return checks


def write_manifest(validation):
    now = dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    lines = [
        "---",
        "stage: deck_build",
        "stage_status: draft",
        "requires_confirmed:",
        "  - ppt_production_brief",
        "  - fact_ledger",
        "  - defense_narrative",
        "  - storyboard",
        "  - speaker_notes_rehearsal",
        "  - defense_qa_backup",
        "  - asset_layout_plan",
        "  - academic_figure_prompt_when_required",
        "  - content_fidelity_qa",
        "allowed_next_stage: ppt-render-qa-loop",
        "confirmed_by:",
        "---",
        "",
        "# PPT Deck Build Manifest v2 Template",
        "",
        "Generated at: `%s`" % now,
        "",
        "## Outputs",
        "",
        "- PPTX: `%s`" % rel(PPTX_OUT),
        "- Build script: `%s`" % rel(os.path.join(ROOT, "exp", "build_marx_report_deck_v2_template.py")),
        "- Manifest: `%s`" % rel(MANIFEST_OUT),
        "- Source template: `%s`" % rel(TEMPLATE),
        "",
        "## Template Background Selection",
        "",
        "| Deck slide | Template slide used | Reason |",
        "|---:|---:|---|",
        "| 1 | 11 | 清华模板正文页，适合方法闭环内容重排 |",
        "| 2 | 13 | 图文/双区块模板页，适合案例材料矩阵 |",
        "| 3 | 23 | 中心图+四周说明模板页，适合五维框架图 |",
        "| 4 | 26 | 时间线/总结类模板页，适合参考文献表格改造 |",
        "",
        "## Change From v1 Image2",
        "",
        "- 直接基于用户提供模板的候选页生成新 PPT，不再只手工套用配色。",
        "- 保留模板的顶部 logo、紫色标题线、淡灰背景纹理和正文页风格。",
        "- 删除模板占位文本和图片占位内容，重排本项目实际内容。",
        "- Slide 3 继续使用已确认的 image2 五维框架图。",
        "",
        "## Counts",
        "",
        "- Slide count: 4",
        "- Main slide count: 4",
        "- Backup slide count: 0",
        "- Notes slides: %s" % count_notes(PPTX_OUT),
        "",
        "## Generated Image Provenance",
        "",
        "- Slide 3 image: `%s`" % rel(IMAGE2),
        "- Prompt artifact: `align/academic_figure_prompt_v0.md`, status confirmed.",
        "- Boundary: explanatory framework image, not evidence image.",
        "",
        "## Structural Validation",
        "",
    ]
    for key in sorted(validation.keys()):
        lines.append("- %s: %s" % (key, validation[key]))
    lines.extend([
        "",
        "## Known Pre-render Risks",
        "",
        "- 使用模板页后整体风格更统一，但仍需 PowerPoint COM render QA 检查是否有占位文本残留、字体替换、表格可读性和图片压边。",
        "- Slide 3 中部图为 raster image，不可逐元素编辑；标题、行动标题、页脚与备注仍为可编辑文本。",
        "",
        "## Required Handoff",
        "",
        "- This is a draft deck build. User confirmation is required before render QA.",
    ])
    with open(MANIFEST_OUT, "w", encoding="utf-8", newline="\n") as f:
        f.write("\n".join(lines) + "\n")


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = build_deck()
    prs.save(PPTX_OUT)
    validation = validate(PPTX_OUT)
    write_manifest(validation)
    print("PPTX:", PPTX_OUT)
    print("Manifest:", MANIFEST_OUT)
    print("Validation:", validation)


if __name__ == "__main__":
    main()
