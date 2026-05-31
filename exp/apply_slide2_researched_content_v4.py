from __future__ import print_function

import os
import zipfile

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
SRC = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v3_template_align.pptx")
OUT = os.path.join(ROOT, "generated_pptx_test", "marx_report_7_8_section_v4_researched_slide2.pptx")


TEXT = {
    0: "案例基础：方法与材料如何进入第八部分",
    3: "第八部分应从“企业介绍/参观感想”改为“多源材料互证后的初步案例分析”",
    4: "材料入口",
    5: "深圳歌尔泰克走访材料\n企业参观、展厅讲解\n座谈交流与实践推文",
    7: "公开补强",
    8: "歌尔股份 2025 年报/摘要\n业务范围与研发背景\n声光电、智能硬件、全球布局",
    10: "分析产出",
    11: "提炼五维框架\n研发设计、工程转化\n产业链协同、人才组织、全球布局",
    12: "改写口径：以深圳歌尔泰克走访材料为入口，结合歌尔股份公开资料补充背景；集团公开数据不直接等同于单一走访主体。",
    13: "来源：实践推文；歌尔股份 2025 年报/摘要；周文、许凌云 2023；尹西明等 2024；Yin 2018",
}

NOTES = """这里我重新调整了第二页的口径。第二页不再只是说明“有哪些材料”，而是回答这些材料如何进入第八部分的案例分析。

左边是材料入口，也就是我们实际掌握的深圳歌尔泰克走访、展厅参观、座谈交流和实践推文；中间是公开补强，用歌尔股份 2025 年报和摘要补充业务范围、研发背景、声光电与智能硬件布局等相对稳定的信息；右边是分析产出，把第八部分收束到研发设计、工程转化、产业链协同、人才组织和全球布局五个维度。

这样写的好处是，案例分析既不只是参观感想，也不变成企业宣传，而是变成有材料来源、有公开资料补强、有理论解释框架的初步分析。同时要注意口径：公开数据属于歌尔股份及子公司层面的资料，不能直接等同于深圳歌尔泰克单一走访主体。"""


def count_notes(path):
    with zipfile.ZipFile(path, "r") as zf:
        return len([n for n in zf.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")])


def set_text(shape, text, font_size=None, bold=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, color=(20, 28, 42)):
    shape.text = text
    tf = shape.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for p in tf.paragraphs:
        p.alignment = align
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.color.rgb = RGBColor(*color)
            if font_size is not None:
                r.font.size = Pt(font_size)
            if bold is not None:
                r.font.bold = bold


def update_notes(slide):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = NOTES


def main():
    prs = Presentation(SRC)
    slide = prs.slides[1]

    for idx, shape in enumerate(slide.shapes):
        if idx not in TEXT or not getattr(shape, "has_text_frame", False):
            continue
        if idx == 0:
            set_text(shape, TEXT[idx], font_size=22, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        elif idx == 3:
            set_text(shape, TEXT[idx], font_size=17, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        elif idx in (4, 7, 10):
            set_text(shape, TEXT[idx], font_size=16, bold=True, color=(255, 255, 255))
        elif idx in (5, 8, 11):
            set_text(shape, TEXT[idx], font_size=14.2, bold=False)
        elif idx == 12:
            set_text(shape, TEXT[idx], font_size=11.2, bold=False, color=(55, 55, 70))
        elif idx == 13:
            set_text(shape, TEXT[idx], font_size=8.8, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, color=(32, 32, 32))

    update_notes(slide)
    prs.save(OUT)
    print("pptx:", OUT)
    print("notes_slides:", count_notes(OUT))


if __name__ == "__main__":
    main()
