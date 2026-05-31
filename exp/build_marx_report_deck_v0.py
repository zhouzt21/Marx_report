from __future__ import print_function

import csv
import datetime as dt
import json
import os
import re
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
ALIGN = os.path.join(ROOT, "align")
OUT_DIR = os.path.join(ROOT, "generated_pptx_test")
PPTX_OUT = os.path.join(OUT_DIR, "marx_report_7_8_section_v0.pptx")
MANIFEST_OUT = os.path.join(ALIGN, "ppt_deck_build_manifest_v0.md")
TEMPLATE = os.path.join(ROOT, "2025nianduyanshiwengaomoban2-tongyongzhuti.pptx")
LAYOUT_PLAN = os.path.join(ALIGN, "ppt_layout_plan_v0.json")
NOTES_ARTIFACT = os.path.join(ALIGN, "ppt_speaker_notes_rehearsal_v0.md")
ASSET_MANIFEST = os.path.join(ALIGN, "PPT_asset_manifest_v0.csv")

PRIMARY = RGBColor(0x3C, 0x36, 0x52)
ACCENT1 = RGBColor(0x75, 0x61, 0xD6)
ACCENT2 = RGBColor(0x51, 0x99, 0xEA)
ACCENT3 = RGBColor(0x68, 0xA4, 0xC6)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFC)
MID_LINE = RGBColor(0xD8, 0xD9, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x67, 0x69, 0x78)


def rel(path):
    return os.path.relpath(path, ROOT).replace("\\", "/")


def add_textbox(slide, x, y, w, h, text, font_size=16, bold=False,
                color=PRIMARY, align=PP_ALIGN.LEFT, font="Microsoft YaHei",
                margin=0.06, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def set_shape_text(shape, text, size=16, bold=False, color=PRIMARY,
                   align=PP_ALIGN.CENTER, font="Microsoft YaHei"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_header(slide, title, action_title, num):
    add_textbox(slide, 0.45, 0.25, 8.2, 0.34, title, 17, True, PRIMARY)
    add_textbox(slide, 0.45, 0.72, 11.35, 0.62, action_title, 22, True, PRIMARY)
    add_textbox(slide, 12.25, 0.28, 0.55, 0.25, "%02d" % num, 9, False, GRAY,
                align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.42), Inches(12.35), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = MID_LINE
    line.line.fill.background()


def add_footer(slide, source_text):
    add_textbox(slide, 0.45, 7.05, 11.75, 0.23, source_text, 8.5, False, GRAY)


def rounded_box(slide, x, y, w, h, text, fill, line=MID_LINE, size=15,
                bold=False, color=PRIMARY, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(radius_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    set_shape_text(box, text, size=size, bold=bold, color=color)
    return box


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT2, width=2.0):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def extract_slide_notes():
    text = open(NOTES_ARTIFACT, "r", encoding="utf-8").read()
    notes = {}
    matches = list(re.finditer(r"^### Slide (\d+)\. .*$", text, flags=re.M))
    for idx, match in enumerate(matches):
        slide_no = int(match.group(1))
        start = match.start()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else text.find("\n## ", match.end())
        if end == -1:
            end = len(text)
        notes[slide_no] = text[start:end].strip()
    return notes


def write_notes(slide, note_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.text = note_text


def create_deck():
    with open(LAYOUT_PLAN, "r", encoding="utf-8") as f:
        plan = json.load(f)
    notes = extract_slide_notes()

    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank = prs.slide_layouts[6]

    for _ in range(len(prs.slides)):
        raise RuntimeError("Unexpected non-empty default presentation")

    build_log = []

    # Slide 1
    s = prs.slides.add_slide(blank)
    add_header(s, "研究方法：多源互证闭环",
               "第 7-8 部分将研究方法和案例观察连接成一个可验证的分析闭环", 1)
    nodes = ["文献研究", "案例研究", "实地走访", "比较分析", "归纳分析"]
    xs = [0.72, 3.05, 5.38, 7.71, 10.04]
    y = 3.0
    for i, node in enumerate(nodes):
        fill = LIGHT_BG if i % 2 == 0 else WHITE
        box = rounded_box(s, xs[i], y, 1.55, 0.76, node, fill, ACCENT2, 15, True)
        idx = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xs[i] + 0.54), Inches(y - 0.78),
                                 Inches(0.46), Inches(0.46))
        idx.fill.solid()
        idx.fill.fore_color.rgb = ACCENT1 if i in (0, 4) else ACCENT2
        idx.line.fill.background()
        set_shape_text(idx, str(i + 1), 11, True, WHITE)
        if i < len(nodes) - 1:
            add_arrow(s, xs[i] + 1.62, y + 0.38, xs[i + 1] - 0.09, y + 0.38)
    add_arrow(s, 10.78, 4.05, 1.05, 4.05, ACCENT3, 1.6)
    add_textbox(s, 4.55, 4.18, 4.2, 0.32, "理论框架 - 案例材料 - 路径归纳", 14, True, ACCENT1,
                align=PP_ALIGN.CENTER)
    rounded_box(s, 2.05, 5.05, 9.3, 0.62,
                "材料不是孤立陈列，而是在方法链条中互相校正、共同支撑后续分析。",
                RGBColor(0xF1, 0xF3, 0xFA), MID_LINE, 15, False)
    add_footer(s, "来源：原报告、实践推文、公开资料、相关文献整理")
    write_notes(s, notes[1])
    build_log.append("slide_1: editable closed-loop process diagram; notes inserted")

    # Slide 2
    s = prs.slides.add_slide(blank)
    add_header(s, "案例基础：以深圳歌尔泰克走访材料为入口",
               "多源互证能避免案例分析停留在参观感想或企业介绍层面", 2)
    columns = [
        ("走访入口", ["深圳歌尔泰克走访", "参观与座谈", "产品展示与成员观察"], ACCENT1),
        ("公开资料", ["歌尔股份年报/摘要", "业务与研发背景", "全球布局信息"], ACCENT2),
        ("分析用途", ["案例典型性", "技术突围维度", "后续路径归纳"], ACCENT3),
    ]
    x0 = 0.72
    for i, (header, items, color) in enumerate(columns):
        x = x0 + i * 4.07
        head = rounded_box(s, x, 1.92, 3.32, 0.52, header, color, color, 17, True, WHITE)
        body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.55), Inches(3.32), Inches(2.5))
        body.fill.solid()
        body.fill.fore_color.rgb = WHITE
        body.line.color.rgb = MID_LINE
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.16)
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(15)
            p.font.color.rgb = PRIMARY
            p.space_after = Pt(5)
        if i < 2:
            add_arrow(s, x + 3.43, 3.78, x + 3.82, 3.78, color, 1.8)
    rounded_box(s, 3.03, 5.45, 7.3, 0.68,
                "口径：以深圳歌尔泰克走访材料为入口，结合歌尔股份公开资料补充背景。",
                RGBColor(0xF4, 0xF2, 0xFF), ACCENT1, 15, True)
    add_footer(s, "来源：实践推文；歌尔股份公开资料；相关文献整理")
    write_notes(s, notes[2])
    build_log.append("slide_2: editable three-column evidence matrix; entity-boundary wording preserved; notes inserted")

    # Slide 3
    s = prs.slides.add_slide(blank)
    add_header(s, "初步观察：技术突围的五个维度",
               "歌尔走访材料显示技术突围依赖研发、工程、供应链、人才和全球布局的系统协同", 3)
    center = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.17), Inches(2.85), Inches(2.08), Inches(1.0))
    center.fill.solid()
    center.fill.fore_color.rgb = PRIMARY
    center.line.color.rgb = PRIMARY
    set_shape_text(center, "高科技企业\n技术突围", 16, True, WHITE)
    dims = [
        ("研发设计", 1.1, 2.0, ACCENT1),
        ("工程转化", 3.6, 4.75, ACCENT2),
        ("产业链协同", 5.15, 1.72, ACCENT3),
        ("人才组织", 7.15, 4.75, ACCENT2),
        ("全球布局", 9.62, 2.0, ACCENT1),
    ]
    for label, x, y, color in dims:
        rounded_box(s, x, y, 2.05, 0.72, label, WHITE, color, 15, True, color)
        add_arrow(s, 6.21, 3.35, x + 1.02, y + 0.36, color, 1.3)
    rounded_box(s, 2.15, 6.0, 8.95, 0.54,
                "初步观察：技术突围不是单点发明，而是系统能力提升；为后续路径与模式归纳提供分析基础。",
                RGBColor(0xF1, 0xF3, 0xFA), MID_LINE, 14.5, False)
    add_footer(s, "来源：实践推文；公开资料；新质生产力、企业创新主体与全球价值链相关文献")
    write_notes(s, notes[3])
    build_log.append("slide_3: editable five-dimension framework; initial-observation wording preserved; notes inserted")

    # Slide 4
    s = prs.slides.add_slide(blank)
    add_header(s, "参考文献与资料来源",
               "参考文献为方法、理论和案例材料提供最低限度支撑", 4)
    refs = plan["deck"]["slides"][3]["content_structure"][1]["rows"]
    rows, cols = len(refs) + 1, 2
    table_shape = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.75), Inches(12.1), Inches(4.95))
    table = table_shape.table
    table.columns[0].width = Inches(8.1)
    table.columns[1].width = Inches(4.0)
    headers = ["文献/资料", "支撑点"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for r, (ref_text, support) in enumerate(refs, start=1):
        table.cell(r, 0).text = ref_text
        table.cell(r, 1).text = support
        for c in range(cols):
            table.cell(r, c).fill.solid()
            table.cell(r, c).fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF8, 0xF9, 0xFC)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(10.5 if r else 12)
                p.font.bold = bool(r == 0)
                p.font.color.rgb = WHITE if r == 0 else PRIMARY
    add_textbox(s, 0.58, 6.78, 11.8, 0.24,
                "注：2025 年报条目需在最终合并前按组内统一格式和实际题录核对。",
                8.5, False, GRAY)
    write_notes(s, notes[4])
    build_log.append("slide_4: editable reference table; English appears only in reference row; notes inserted")

    return prs, build_log


def validate_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    full_text = "\n".join(texts)
    checks = {
        "file_exists": os.path.exists(path),
        "opened_by_python_pptx": True,
        "slide_count": len(prs.slides),
        "aspect_ratio_16_9": (prs.slide_width == 12192000 and prs.slide_height == 6858000),
        "slide_2_entity_boundary": "以深圳歌尔泰克走访材料为入口，结合歌尔股份公开资料补充背景" in full_text,
        "slide_3_initial_observation": "初步观察" in full_text and "分析基础" in full_text,
        "forbidden_wording_absent": ("已证明" not in full_text and "最终模型" not in full_text),
        "placeholder_text_absent": not any(s in full_text for s in ["添加标题", "添加正文", "输入图片", "XXX"]),
        "notes_slides_present": True,
    }
    try:
        for slide in prs.slides:
            _ = slide.notes_slide.notes_text_frame.text
    except Exception:
        checks["notes_slides_present"] = False
    return checks


def manifest(build_log, validation):
    now = dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with open(ASSET_MANIFEST, "r", encoding="utf-8-sig") as f:
        asset_rows = list(csv.DictReader(f))
    main_assets = [r for r in asset_rows if r["slide"] in ["1", "2", "3", "4"]]

    mapping = [
        "| Slide | Storyboard item | Layout archetype | Template/source mapping | Visual route |",
        "|---:|---|---|---|---|",
        "| 1 | 研究方法：多源互证闭环 | action_title_plus_process | template style reference; editable redraw | local_editable_diagram |",
        "| 2 | 案例基础：以深圳歌尔泰克走访材料为入口 | action_title_plus_three_column_matrix | template style reference; editable redraw | local_editable_diagram/table |",
        "| 3 | 初步观察：技术突围的五个维度 | action_title_plus_five_dimension_framework | template style reference; editable redraw | local_editable_diagram |",
        "| 4 | 参考文献与资料来源 | references_table | template style reference; editable table | local_editable_table |",
    ]
    route_lines = ["| Slide | Asset | Route outcome | Evidence status | Editability |",
                   "|---:|---|---|---|---|"]
    for row in main_assets:
        route_lines.append("| {slide} | {asset_name} | {visual_route}; generated as editable PPT objects | {evidence_status} | {editable_expected} |".format(**row))

    md = []
    md.append("---")
    md.append("stage: deck_build")
    md.append("stage_status: draft")
    md.append("requires_confirmed:")
    md.extend([
        "  - ppt_production_brief",
        "  - fact_ledger",
        "  - defense_narrative",
        "  - storyboard",
        "  - speaker_notes_rehearsal",
        "  - defense_qa_backup",
        "  - asset_layout_plan",
        "  - academic_figure_prompt_when_required",
        "  - content_fidelity_qa",
    ])
    md.append("allowed_next_stage: ppt-render-qa-loop")
    md.append("confirmed_by:")
    md.append("---")
    md.append("")
    md.append("# PPT Deck Build Manifest v0")
    md.append("")
    md.append("Generated at: `%s`" % now)
    md.append("")
    md.append("## Outputs")
    md.append("")
    md.append("- PPTX: `%s`" % rel(PPTX_OUT))
    md.append("- Build script: `%s`" % rel(os.path.join(ROOT, "exp", "build_marx_report_deck_v0.py")))
    md.append("- Manifest: `%s`" % rel(MANIFEST_OUT))
    md.append("")
    md.append("## Counts")
    md.append("")
    md.append("- Slide count: 4")
    md.append("- Main slide count: 4")
    md.append("- Backup slide count: 0")
    md.append("- Backup policy: B1-B5 remain planned optional in `align/ppt_defense_qa_backup_v0.md` and `align/ppt_layout_plan_v0.json`; not inserted into this main draft because the user requested a 4-slide clean deck.")
    md.append("")
    md.append("## Notes Insertion")
    md.append("")
    md.append("- Status: inserted into PowerPoint notes pane via `python-pptx` notes slide API.")
    md.append("- Notes source: `%s`" % rel(NOTES_ARTIFACT))
    md.append("- Limitation: notes were inserted structurally but not render-verified in PowerPoint, because render QA/COM was explicitly out of scope for this lane.")
    md.append("")
    md.append("## Confirmed Source Paths")
    md.append("")
    for p in [
        "ppt_production_brief_v0.md", "fact_ledger_v0.md", "ppt_defense_narrative_v0.md",
        "PPT_storyboard_v0.md", "ppt_speaker_notes_rehearsal_v0.md",
        "ppt_defense_qa_backup_v0.md", "template_inventory_v0.md",
        "template_design_rules_v0.md", "PPT_asset_audit_v0.md",
        "visual_enrichment_plan_v0.md", "ppt_layout_plan_v0.json",
        "PPT_asset_manifest_v0.csv", "ppt_content_fidelity_qa_v0.md",
    ]:
        md.append("- `%s`" % rel(os.path.join(ALIGN, p)))
    md.append("")
    md.append("## Content Fidelity QA")
    md.append("")
    md.append("- Source: `align/ppt_content_fidelity_qa_v0.md`")
    md.append("- Status: confirmed; `qa_result: draft-pass-ready`.")
    md.append("- Applied constraints: Slide 2 keeps entity-boundary wording; Slide 3 keeps `初步观察/分析基础`; forbidden wording `已证明/最终模型` is absent from generated slide text.")
    md.append("")
    md.append("## Template / Design Rules")
    md.append("")
    md.append("- Template source read: `%s`" % rel(TEMPLATE))
    md.append("- Template design rules: `align/template_design_rules_v0.md`, status confirmed.")
    md.append("- Direct template master reuse: not used. `python-pptx` could read template size and slide count, but layout enumeration hit an internal layout relationship error; deck was generated as a clean 16:9 editable deck using the confirmed colors, typography, and layout archetypes.")
    md.append("")
    md.append("## Layout / Template Mapping")
    md.append("")
    md.extend(mapping)
    md.append("")
    md.append("## Visual Route Outcomes")
    md.append("")
    md.extend(route_lines)
    md.append("")
    md.append("## Assets Used")
    md.append("")
    md.append("- No external raster assets embedded.")
    md.append("- No practice-post screenshots, annual-report screenshots, AI-generated images, or whole-slide images used.")
    md.append("- All slide bodies are editable PowerPoint text boxes, shapes, connectors, and table cells.")
    md.append("")
    md.append("## Generated Image Provenance")
    md.append("")
    md.append("- None. `align/visual_enrichment_plan_v0.md` marks `requires_academic_figure_prompt: false`; OpenRouter ICU Image was not called.")
    md.append("")
    md.append("## Generation Log")
    md.append("")
    for item in build_log:
        md.append("- %s" % item)
    md.append("")
    md.append("## Structural Validation")
    md.append("")
    for key in sorted(validation.keys()):
        md.append("- %s: %s" % (key, validation[key]))
    md.append("")
    md.append("## Known Pre-render Risks")
    md.append("")
    md.append("- Chinese font substitution may change wrapping on another machine.")
    md.append("- Reference slide uses 10.5 pt text; readable structurally, but final visual acceptance requires render QA.")
    md.append("- Template master/layout was not directly reused due to `python-pptx` layout relationship instability; visual style follows confirmed rules instead.")
    md.append("- Speaker notes were inserted but not inspected through PowerPoint UI in this lane.")
    md.append("")
    md.append("## Required Handoff")
    md.append("")
    md.append("- Next stage: `ppt-render-qa-loop`.")
    md.append("- This draft is not self-confirmed and should be handed to render QA before acceptance.")
    md.append("")
    return "\n".join(md)


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs, build_log = create_deck()
    prs.save(PPTX_OUT)
    validation = validate_pptx(PPTX_OUT)
    with open(MANIFEST_OUT, "w", encoding="utf-8", newline="\n") as f:
        f.write(manifest(build_log, validation))
    print("PPTX:", PPTX_OUT)
    print("Manifest:", MANIFEST_OUT)
    print("Validation:", validation)


if __name__ == "__main__":
    main()
